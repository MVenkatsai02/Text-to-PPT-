import streamlit as st
import base64
import google.generativeai as genai
import pptx
from pptx.util import Pt
import os
import re
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv('gemini_api_key'))

# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)
MAX_SLIDE_CONTENT_LENGTH = 400  # Max characters per slide before creating a new one

def generate_slide_titles(topic):
    """Generate subtopics for the presentation."""
    prompt = f"Generate 10 detailed subtopics for the topic '{topic}'. Each subtopic should be clear and concise, without numbering or extra symbols."
    model = genai.GenerativeModel("gemini-1.5-flash")
    response = model.generate_content(prompt)
    return [title.strip() for title in response.text.strip().split("\n") if title.strip()]

def generate_slide_content(slide_title):
    """Generate content for a given slide title."""
    prompt = f"Generate detailed bullet points for the slide '{slide_title}'. Each point should have a bold title followed by a brief explanation. Avoid using asterisks or extra symbols."
    model = genai.GenerativeModel("gemini-1.5-flash")
    response = model.generate_content(prompt)
    return response.text.strip()

def clean_text(text):
    """Remove unwanted formatting like asterisks from bold text."""
    return re.sub(r'\*\*(.*?)\*\*', r'\1', text)

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

MAX_BULLETS_PER_SLIDE = 6  # Define the maximum number of bullet points per slide

def split_content_into_slides(content):
    """Split content into slides while maximizing space utilization."""
    paragraphs = content.split("\n")
    slides = []
    current_slide = []
    
    for para in paragraphs:
        para = clean_text(para)  # Clean text from unwanted symbols
        if len(current_slide) < MAX_BULLETS_PER_SLIDE:
            current_slide.append(para)
        else:
            slides.append("\n".join(current_slide))
            current_slide = [para]

    if current_slide:
        slides.append("\n".join(current_slide))

    return slides


def create_presentation(topic, slide_titles, slide_contents):
    """Create a PowerPoint presentation while dynamically adjusting content layout."""
    prs = pptx.Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Add title slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = topic.upper()
    title_slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        content_slides = split_content_into_slides(slide_content)
        first_slide = True

        for index, content in enumerate(content_slides):
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = slide_title if first_slide else f"{slide_title} (contd.)"

            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.clear()  # Clear any existing text to format properly

            for bullet in content.split("\n"):
                if bullet.strip():
                    p = text_frame.add_paragraph()
                    p.text = bullet.strip()
                    p.space_after = Pt(8)
                    p.level = 0  # Ensures it stays as a bullet point
                    p.alignment = PP_ALIGN.LEFT

            slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.font.size = SLIDE_FONT_SIZE

            first_slide = False  # Mark first slide as completed

    output_dir = "generated_ppt"
    os.makedirs(output_dir, exist_ok=True)
    ppt_filename = os.path.join(output_dir, f"{topic.replace(' ', '_')}_presentation.pptx")
    prs.save(ppt_filename)
    return ppt_filename

def get_ppt_download_link(ppt_filename):
    """Generate a download link for the PowerPoint presentation."""
    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()
    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{os.path.basename(ppt_filename)}">Download the PowerPoint Presentation</a>'

def main():
    """Streamlit UI for the AI-powered PowerPoint Generator."""
    st.title("AI-Powered PowerPoint Generator")
    topic = st.text_input("Enter the topic for your presentation:")
    generate_button = st.button("Generate Presentation")
    
    if generate_button and topic:
        st.info("Generating presentation... Please wait.")
        slide_titles = generate_slide_titles(topic)
        slide_contents = [generate_slide_content(title) for title in slide_titles]
        ppt_filename = create_presentation(topic, slide_titles, slide_contents)
        st.success("Presentation generated successfully!")
        st.markdown(get_ppt_download_link(ppt_filename), unsafe_allow_html=True)

if __name__ == "__main__":
    main()
